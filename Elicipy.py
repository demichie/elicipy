import datetime
import numpy as np
import os
import sys
import pandas as pd

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx import Presentation

from saveFromGithub import saveDataFromGithub
from createSamples import createSamples
from COOKEweights import COOKEweights
from ERFweights import ERFweights
from merge_csv import merge_csv

# from createPlots import create_fig_hist
from createPlots import create_figure_violin
from createPlots import create_figure_pie
from createPlots import create_figure_trend
from createPlots import create_figure_answers
from createPlots import create_barplot
from createPlots import create_figure_index

from tools import printProgressBar

# from krippendorff_alpha import calculate_alpha
from computeIndex import calculate_index

max_len_table = 21
max_len_tableB = 18

max_len_plot = 21


def add_date(slide):

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.2),
                                   Inches(16), Inches(0.3))
    shape.shadow.inherit = False
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(60, 90, 180)
    line = shape.line
    line.color.rgb = RGBColor(60, 90, 180)
    shape.text = "Expert elicitation " + \
        datetime.datetime.today().strftime("%d-%b-%Y")
    shape_para = shape.text_frame.paragraphs[0]
    shape_para.font.name = "Helvetica"
    shape_para.font.size = Pt(17)


def add_small_logo(slide, left, top, logofile):

    slide.shapes.add_picture(logofile,
                             left + Inches(13.0),
                             top + Inches(6.8),
                             width=Inches(0.8))


def add_figure(slide, figname, left, top, width):

    img = slide.shapes.add_picture(figname,
                                   left + Inches(3.4),
                                   top,
                                   width=width)

    slide.shapes._spTree.insert(2, img._element)


def add_small_figure(slide, figname, left, top, width):

    img = slide.shapes.add_picture(figname, left, top, width=width)

    slide.shapes._spTree.insert(2, img._element)


def add_title(slide, text_title):

    title_shape = slide.shapes.title
    title_shape.text = text_title
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(15)
    title_shape.height = Inches(2)
    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.name = "Helvetica"
    if len(text_title) < 50:

        title_para.font.size = Pt(44)

    else:

        title_para.font.size = Pt(34)


def add_text_box(slide, left, top, text_box, font_size):

    txBox = slide.shapes.add_textbox(left - Inches(1),
                                     top + Inches(0.5),
                                     width=Inches(4),
                                     height=Inches(5))
    tf = txBox.text_frame
    tf.text = text_box
    for par in tf.paragraphs:

        par.font.size = Pt(font_size)
        # tf.paragraphs[0].font.size = Pt(font_size)
    # tf.text = 'prova'
    tf.word_wrap = True


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def read_answers(input_dir, csv_file, group, n_pctl, df_indexes_SQ,
                 df_indexes_TQ, seed, target, output_dir, elicitation_name,
                 write_flag, label_indexes):

    verbose = False

    try:

        from ElicipyDict import label_flag

    except ImportError:

        label_flag = False

    import difflib

    # merge the files of the different experts
    # creating one file for seeds and one for tagets
    merge_csv(input_dir, seed, target, group, csv_file, label_flag, write_flag)

    if seed:

        # seeds file name
        filename = input_dir + "/seed.csv"

        # Read a comma-separated values (csv) file into DataFrame df_SQ
        df_SQ = pd.read_csv(filename)

        # create a 2D numpy array with the answers to the seed questions
        cols_as_np = df_SQ[df_SQ.columns[4:]].to_numpy()

        # we want to work with a 3D array, with the following dimension:
        # n_expert X n_pctl X n_SQ
        n_experts = cols_as_np.shape[0]
        n_SQ = int(cols_as_np.shape[1] / n_pctl)

        # The second column (index 1) must contain the first name of the expert
        firstname = df_SQ[df_SQ.columns[1]].astype(str).tolist()

        # The third column (index 2) must contain the first name of the expert
        surname = df_SQ[df_SQ.columns[2]].astype(str).tolist()

        # create a list with firstname+surname
        # this is needed to search for expert matches between
        # seed and target questions
        NS_SQ = []

        for name, surname in zip(firstname, surname):

            NS_SQ.append(name + surname)

        if verbose:
            print("NS_SQ", NS_SQ)

        # reshaped numpy array with expert answers
        SQ_array = np.reshape(cols_as_np, (n_experts, n_SQ, n_pctl))

        # swap the array to have seed for the last index
        SQ_array = np.swapaxes(SQ_array, 1, 2)

        # sort according to the percentile values
        # (sometimes the expert give the percentiles in the wrong order)
        SQ_array = np.sort(SQ_array, axis=1)

        # chech and correct for equal percentiles
        for i in np.arange(n_SQ):

            for k in np.arange(n_experts):

                # if 5% and 50% percentiles are equal, reduce 5%
                if SQ_array[k, 0, i] == SQ_array[k, 1, i]:

                    SQ_array[k, 0, i] = SQ_array[k, 1, i] * 0.99

                # if 50% and 95% percentiles are equal, increase 95%
                if SQ_array[k, 2, i] == SQ_array[k, 1, i]:

                    SQ_array[k, 2, i] = SQ_array[k, 1, i] * 1.01

        # if we have a subset of the SQ, then extract from SQ_array
        # the correct slice
        if len(df_indexes_SQ) > 0:

            # print('SQ_array',SQ_array)
            SQ_array = SQ_array[:, :, df_indexes_SQ]
            n_SQ = len(df_indexes_SQ)
            # print('SQ_array',SQ_array)

        if verbose:

            for i in np.arange(n_SQ):

                print("")
                print("Seed question ", label_indexes[i])
                print(SQ_array[:, :, i])

    else:

        n_SQ = 0
        df_indexes_SQ = []

    if target:

        filename = input_dir + "/target.csv"

        # Read a comma-separated values (csv) file into DataFrame df_TQ
        df_TQ = pd.read_csv(filename)

        # The second column (index 1) must contain the first name of the expert
        firstname = df_TQ[df_TQ.columns[1]].astype(str).tolist()

        # The third column (index 2) must contain the first name of the expert
        surname = df_TQ[df_TQ.columns[2]].astype(str).tolist()

        # create a list with firstname+surname
        # this is needed to search for expert matches between seed
        # and target questions
        NS_TQ = []

        for name, surname in zip(firstname, surname):

            NS_TQ.append(name + surname)

        if seed:

            sorted_idx = []

            # loop to search for matches between experts in seed and target
            for SQ_name in NS_SQ:

                index = NS_TQ.index(
                    difflib.get_close_matches(SQ_name, NS_TQ)[0])
                sorted_idx.append(index)

            if verbose:
                print("Sorted list of experts to match the order of seeds:",
                      sorted_idx)

                print(NS_SQ)
                print([NS_TQ[s_idx] for s_idx in sorted_idx])

        else:

            sorted_idx = range(len(NS_TQ))
            NS_SQ = NS_TQ

        # create a 2D numpy array with the answers to the target questions
        cols_as_np = df_TQ[df_TQ.columns[4:]].to_numpy()

        # sort for expert names
        cols_as_np = cols_as_np[sorted_idx, :]

        # we want to work with a 3D array, with the following dimension:
        # n_expert X n_pctl X n_TQ
        n_experts = cols_as_np.shape[0]

        n_TQ = int(cols_as_np.shape[1] / n_pctl)

        # reshaped numpy array with expert answers
        TQ_array = np.reshape(cols_as_np, (n_experts, n_TQ, n_pctl))

        # swap the array to have pctls for the second index
        TQ_array = np.swapaxes(TQ_array, 1, 2)

        # sort according to the percentile values
        # (sometimes the expert give the percentiles in the wrong order)
        TQ_array = np.sort(TQ_array, axis=1)

        for i in np.arange(n_TQ):

            for k in np.arange(n_experts):
                if TQ_array[k, 0, i] == TQ_array[k, 1, i]:
                    TQ_array[k, 0, i] = TQ_array[k, 1, i] * 0.99
                if TQ_array[k, 2, i] == TQ_array[k, 1, i]:
                    TQ_array[k, 2, i] = TQ_array[k, 1, i] * 1.01

        if len(df_indexes_TQ) > 0:

            # print('TQ_array',TQ_array)
            TQ_array = TQ_array[:, :, df_indexes_TQ]
            n_TQ = len(df_indexes_TQ)
            # print('TQ_array',TQ_array)

        if verbose:

            for i in np.arange(n_TQ):

                print("Target question ", label_indexes[i + n_SQ])
                print(TQ_array[:, :, i])

    else:

        n_TQ = 0
        TQ_array = np.zeros((n_experts, n_pctl, n_TQ))

    if not seed:

        SQ_array = np.zeros((n_experts, n_pctl, n_SQ))

    csv_name = output_dir + "/" + elicitation_name + "_experts.csv"

    d = {"index": range(1, n_experts + 1), "Expert": NS_SQ}

    df = pd.DataFrame(data=d)

    df.to_csv(csv_name, index=False)

    return n_experts, n_SQ, n_TQ, SQ_array, TQ_array, NS_SQ


def read_questionnaire(input_dir, csv_file, seed, target):
    """Read .csv questionnaire file

    Parameters
    ----------
    input_dir : string
        name of input folder
    csv_file : string
        name of csv_file
    target : boolean
        True to read targets

    Returns
    -------
    df_indexes_SQ, df_indexes_TQ, SQ_scale, SQ_realization, TQ_scale, \
        SQ_minVals, SQ_maxVals, TQ_minVals, TQ_maxVals, SQ_units, TQ_units, \
        SQ_LongQuestion, TQ_LongQuestion, SQ_question, TQ_question, \
        idx_list, global_scale, global_log, label_indexes, parents, \
        global_idxMin, global_idxMax, global_sum50

    """

    verbose = False

    print("STEP1: Reading questionnaire")

    try:

        from ElicipyDict import label_flag

    except ImportError:

        label_flag = False

    df_read = pd.read_csv(input_dir + "/" + csv_file, header=0)
    # print(df_read)

    quest_type = df_read["QUEST_TYPE"].to_list()
    n_SQ = quest_type.count("seed")

    if verbose:
        print(quest_type)

    try:

        from ElicipyDict import seed_list

        if verbose:
            print("seed_list read", seed_list)

    except ImportError:

        seed_list = list(df_read["IDX"])[0:n_SQ]

    if verbose:
        print("seed_list", seed_list)

    # extract the seed questions with index in
    # seed_list (from column IDX)
    new_indexes = []

    for idx in seed_list[:]:

        indices = df_read.index[(df_read["IDX"] == idx)
                                & (df_read.QUEST_TYPE.str.contains("seed"))]

        new_indexes.append(indices[0])

    df_SQ = df_read.iloc[new_indexes]
    df_indexes_SQ = np.array(new_indexes).astype(int)

    if verbose:
        print("df_indexes_SQ", df_indexes_SQ)

    if target:

        try:

            from ElicipyDict import target_list
            if verbose:
                print("target_list read", target_list)

        except ImportError:

            if verbose:
                print("ImportError")
            target_list = list(df_read["IDX"])[n_SQ:]

        if verbose:
            print("target_list", target_list)

        # extract the target questions with index in
        # target_list (from column IDX)
        new_indexes = []

        for idx in target_list[:]:

            indices = df_read.index[
                (df_read["IDX"] == idx)
                & (df_read.QUEST_TYPE.str.contains("target"))]

            new_indexes.append(indices[0])

        df_TQ = df_read.iloc[new_indexes]
        df_indexes_TQ = np.array(new_indexes).astype(int) - n_SQ

        if verbose:
            print("df_indexes_TQ", df_indexes_TQ)

        if seed:

            # df_quest = df_SQ.append(df_TQ)
            df_quest = pd.concat([df_SQ, df_TQ])

        else:

            df_quest = df_TQ

    else:

        df_indexes_TQ = []
        df_quest = df_SQ

    if label_flag:
        label_indexes = df_quest["LABEL"].astype(str).tolist()
    else:
        label_indexes = np.asarray(df_quest["IDX"])
        label_indexes = label_indexes.astype(str).tolist()

    if verbose:
        print("label_indexes", label_indexes)

    data_top = df_quest.head()

    langs = []

    # check if there are multiple languages
    for head in data_top:

        if "LONG Q" in head:

            string = head.replace("LONG Q", "")
            string2 = string.replace("_", "")

            langs.append(string2)

    if verbose:
        print("Languages:", langs)

    try:

        from ElicipyDict import language

    except ImportError:

        language = ""

    # select the columns to use according with the language
    if len(langs) > 1:

        if language in langs:

            lang_index = langs.index(language)
            # list of column indexes to use
            index_list = [1, 2, 3, lang_index + 4] + list(
                range(len(langs) + 4,
                      len(langs) + 15))

        else:

            raise Exception("Sorry, language is not in questionnaire")

    else:

        lang_index = 0
        language = ""
        index_list = list(range(1, 16))

    # list with the short title of the target questions
    SQ_question = []
    # list with the long title of the target questions
    SQ_LongQuestion = []
    # list with min vals for target questions
    SQ_minVals = []
    # list with max vals for target questions
    SQ_maxVals = []
    # list with the units of the target questions
    SQ_units = []
    # scale for target question:
    SQ_scale = []

    SQ_realization = []

    global_log = []

    global_idxMin = []
    global_idxMax = []
    global_sum50 = []

    for i in df_quest.itertuples():

        (
            idx,
            label,
            shortQ,
            longQ,
            unit,
            scale,
            minVal,
            maxVal,
            realization,
            question,
            idxMin,
            idxMax,
            sum50,
            parent,
            image,
        ) = [i[j] for j in index_list]

        minVal = float(minVal)
        maxVal = float(maxVal)

        if scale == "uni":

            global_log.append(0)

        else:

            global_log.append(1)

        if question == "seed":

            SQ_question.append(shortQ)
            SQ_LongQuestion.append(longQ)
            SQ_units.append(unit)
            SQ_scale.append(scale)

            SQ_realization.append(float(realization))

            if minVal.is_integer():

                minVal = int(minVal)

            if maxVal.is_integer():

                maxVal = int(maxVal)

            SQ_minVals.append(minVal)
            SQ_maxVals.append(maxVal)

            if idxMin > 0:
                df_min = df_quest.loc[(df_quest['IDX'] == idxMin)
                                      & (df_quest['QUEST_TYPE'] == 'seed')]
                global_idxMin.append(str(df_min['IDX'].iloc[0]))
            else:
                global_idxMin.append('')

            if idxMax > 0:
                df_max = df_quest.loc[(df_quest['IDX'] == idxMax)
                                      & (df_quest['QUEST_TYPE'] == 'seed')]
                global_idxMax.append(str(df_max['IDX'].iloc[0]))
            else:
                global_idxMax.append('')

            global_sum50.append(sum50)

    # print on screen the units
    if verbose:
        print("Seed_units = ", SQ_units)

    # print on screen the units
    if verbose:
        print("Seed_scales = ", SQ_scale)

    # list with the "title" of the target questions
    TQ_question = []
    # list with the long title of the target questions
    TQ_LongQuestion = []
    TQ_minVals = []
    TQ_maxVals = []
    # list with the units of the target questions
    TQ_units = []
    # scale for target question:
    TQ_scale = []

    idx_list = []
    parents = []

    if target:

        for i in df_quest.itertuples():

            (
                idx,
                label,
                shortQ,
                longQ,
                unit,
                scale,
                minVal,
                maxVal,
                realization,
                question,
                idxMin,
                idxMax,
                sum50,
                parent,
                image,
            ) = [i[j] for j in index_list]

            minVal = float(minVal)
            maxVal = float(maxVal)
            if question == "target":

                TQ_question.append(shortQ)
                TQ_LongQuestion.append(longQ)
                TQ_units.append(unit)
                TQ_scale.append(scale)

                if minVal.is_integer():

                    minVal = int(minVal)

                if maxVal.is_integer():

                    maxVal = int(maxVal)

                TQ_minVals.append(minVal)
                TQ_maxVals.append(maxVal)
                idx_list.append(int(idx))
                parents.append(int(parent))

                if (idxMin > 0):
                    df_min = df_quest.loc[(df_quest['IDX'] == idxMin) &
                                          (df_quest['QUEST_TYPE'] == 'target')]
                    global_idxMin.append(str(df_min['IDX'].iloc[0]))
                else:
                    global_idxMin.append('')

                if (idxMax > 0):
                    df_max = df_quest.loc[(df_quest['IDX'] == idxMax) &
                                          (df_quest['QUEST_TYPE'] == 'target')]
                    global_idxMax.append(str(df_max['IDX'].iloc[0]))
                else:
                    global_idxMax.append('')

                global_sum50.append(sum50)

        # print on screen the units
        if verbose:
            print("Target units = ", TQ_units)

        # print on screen the units
        if verbose:
            print("Target scales = ", TQ_scale)

        global_scale = SQ_scale + TQ_scale

    else:

        global_scale = SQ_scale

    return (df_indexes_SQ, df_indexes_TQ, SQ_scale, SQ_realization, TQ_scale,
            SQ_minVals, SQ_maxVals, TQ_minVals, TQ_maxVals, SQ_units, TQ_units,
            SQ_LongQuestion, TQ_LongQuestion, SQ_question, TQ_question,
            idx_list, global_scale, global_log, label_indexes, parents,
            global_idxMin, global_idxMax, global_sum50, df_quest)


def answer_analysis(input_dir, csv_file, n_experts, n_SQ, n_TQ, SQ_array,
                    TQ_array, realization, global_scale, global_log, alpha,
                    overshoot, cal_power, ERF_flag, Cooke_flag, seed,
                    NS_experts, weights_file):

    verbose = False

    if Cooke_flag < 0:

        # when Cooke_flag is negative, the weights are read from
        # and external file (weights_file) define in input

        from merge_csv import similar

        weights = pd.read_csv(weights_file)
        print(weights)

        fname = weights["First Name"].to_list()
        lname = weights["Last Name"].to_list()

        W = np.zeros((n_experts, 5))

        for i, (f, l) in enumerate(zip(fname, lname)):

            flname = str(f) + str(l)
            lfname = str(l) + str(f)

            for ex, name in enumerate(NS_experts):

                sim1 = similar(flname, name)
                sim2 = similar(lfname, name)
                sim = max(sim1, sim2)

                if sim > 0.8:

                    W[ex, 0] = weights.C[i]
                    W[ex, 1] = weights.I_tot[i]
                    W[ex, 2] = weights.I_real[i]
                    W[ex, 3] = weights.w[i]
                    W[ex, 4] = weights.normW[i]

        W[:, 4] /= np.sum(W[:, 4])

    elif Cooke_flag > 0:

        W, score, information, M = COOKEweights(SQ_array, TQ_array,
                                                realization, alpha,
                                                global_scale, overshoot,
                                                cal_power, Cooke_flag)

        if verbose:

            print('score', score)
            print('information', information)

    else:

        W = np.ones((n_experts, 5))
        score = np.zeros(n_experts)
        information = np.zeros(n_experts)

    if ERF_flag:

        W_erf, score_erf = ERFweights(realization, SQ_array)

    else:

        W_erf = np.ones((n_experts, 5))
        # score_erf = np.zeros(n_experts)

    sensitivity = False

    if sensitivity:

        for i in range(n_SQ):

            SQ_temp = np.delete(SQ_array, i, axis=2)
            realization_temp = np.delete(realization, i)
            global_scale_temp = np.delete(global_scale, i)

            W_temp, score_temp, information_temp, M = COOKEweights(
                SQ_temp, TQ_array, realization_temp, alpha, global_scale_temp,
                overshoot, cal_power, Cooke_flag)

            W_reldiff = W[:, 3] / \
                np.sum(W[:, 3]) - W_temp[:, 3] / np.sum(W_temp[:, 3])

            W_mean = np.mean(np.abs(W_reldiff))
            W_std = np.sqrt(np.sum(np.square(W_reldiff)) / n_experts)

            print(i + 1, W_mean, W_std, np.sum(W_temp[:, 4] > 0))

            W_erf_temp, score_erf_temp = ERFweights(realization_temp, SQ_temp)

            W_reldiff = W_erf[:, 4] - W_erf_temp[:, 4]

            W_mean = np.mean(np.abs(W_reldiff))
            W_std = np.sqrt(np.sum(np.square(W_reldiff)) / n_experts)

            print(i + 1, W_mean, W_std, np.sum(W_erf_temp > alpha))

    Weq = np.ones(n_experts)
    Weqok = [x / n_experts for x in Weq]

    W_gt0_01 = []
    Werf_gt0_01 = []
    expin = []

    k = 1
    for x, y in zip(W[:, 4], W_erf[:, 4]):
        if (x > 0) or (y > 0):
            W_gt0_01.append(x)
            Werf_gt0_01.append(y)
            expin.append(k)
        k += 1

    W_gt0 = [round((x * 100.0), 5) for x in W_gt0_01]
    Werf_gt0 = [round((x * 100.0), 5) for x in Werf_gt0_01]

    if verbose:

        if ERF_flag > 0:

            print("")
            print("W_erf")
            print(W_erf[:, -1])

        if Cooke_flag != 0:

            print("")
            print("W_cooke")
            print(W[:, -1])
            print("score", score)
            print("information", information)
            print("unNormalized weight", W[:, 3])

        print("")
        print("Weq")
        print(Weqok)

    return W, W_erf, Weqok, W_gt0, Werf_gt0, expin


def save_dtt_rll(input_dir, n_experts, n_SQ, n_TQ, df_quest, target,
                 SQ_realization, SQ_scale, SQ_array, TQ_scale, TQ_array):

    # ----------------------------------------- #
    # ------------ Save dtt and rls ----------- #
    # ----------------------------------------- #

    # Save a reference to the original standard output
    original_stdout = sys.stdout

    filename = input_dir + "/seed_and_target.rls"

    with open(filename, "w") as f:
        sys.stdout = f  # Change the standard output to the file we created.

        for i in np.arange(n_SQ):

            # print(i+1,str(i+1),SQ_realization[i],SQ_scale[i])

            print(f'{i+1:>5d} {"SQ"+str(i+1):>13} {""} ' +
                  f'{SQ_realization[i]:6e} {SQ_scale[i]:4}')

        # Reset the standard output to its original value
        sys.stdout = original_stdout

    if target:

        filename = input_dir + "/seed_and_target.dtt"

    else:

        filename = input_dir + "/seed.dtt"

    with open(filename, "w") as f:
        sys.stdout = f  # Change the standard output to the file we created.

        print("* CLASS ASCII OUTPUT FILE. NQ=   3   QU=   5  50  95")
        print("")
        print("")

        for k in np.arange(n_experts):

            for i in np.arange(n_SQ):

                print(f'{k+1:5d} {"Exp"+str(k+1):>8} {i+1:4d} ' +
                      f'{"SQ"+str(i+1):>13} {SQ_scale[i]:4} ' +
                      f'{SQ_array[k, 0, i]:6e} {""} {SQ_array[k, 1, i]:6e} ' +
                      f'{" "}{SQ_array[k, 2, i]:6e}')

            if target:

                for i in np.arange(n_TQ):

                    print(
                        f'{k+1:5d} {"Exp"+str(k+1):>8} {i+1:4d} ' +
                        f'{"TQ"+str(i+1):>13} {TQ_scale[i]:4} ' +
                        f'{TQ_array[k, 0, i]:6e} {""} {TQ_array[k, 1, i]:6e} '
                        + f'{" "}{TQ_array[k, 2, i]:6e}')

        # Reset the standard output to its original value
        sys.stdout = original_stdout

    return


def create_samples(group, n_experts, n_SQ, n_TQ, n_pctl, SQ_array, TQ_array,
                   n_sample, W, W_erf, Weqok, W_gt0, Werf_gt0, expin,
                   global_log, global_minVal, global_maxVal, label_indexes,
                   ERF_flag, Cooke_flag, EW_flag, overshoot, globalSum,
                   normalizeSum):

    verbose = False

    DAT = np.zeros((n_experts * (n_SQ + n_TQ), n_pctl + 2))

    DAT[:, 0] = np.repeat(np.arange(1, n_experts + 1), n_SQ + n_TQ)
    DAT[:, 1] = np.tile(np.arange(1, n_SQ + n_TQ + 1), n_experts)

    DAT[:, 2:] = np.append(SQ_array, TQ_array,
                           axis=2).transpose(0, 2, 1).reshape(-1, 3)
    q_Cooke = np.zeros((n_SQ + n_TQ, 4))
    q_erf = np.zeros((n_SQ + n_TQ, 4))
    q_EW = np.zeros((n_SQ + n_TQ, 4))

    samples = np.zeros((n_sample, n_SQ + n_TQ))
    samples_erf = np.zeros((n_sample, n_SQ + n_TQ))
    samples_EW = np.zeros((n_sample, n_SQ + n_TQ))

    print('       Creating samples for questions')

    for j in np.arange(n_SQ + n_TQ):

        if (n_SQ + n_TQ > 1):

            printProgressBar(j, n_SQ + n_TQ - 1, prefix='      ')

        C_EW = createSamples(
            DAT,
            j,
            Weqok,
            n_sample,
            global_log[j],
            [global_minVal[j], global_maxVal[j]],
            overshoot,
            0,
        )

        samples_EW[:, j] = C_EW

        if Cooke_flag != 0:

            C = createSamples(
                DAT,
                j,
                W[:, 4].flatten(),
                n_sample,
                global_log[j],
                [global_minVal[j], global_maxVal[j]],
                overshoot,
                0,
            )

        else:

            C = C_EW

        samples[:, j] = C

        if ERF_flag > 0:

            C_erf = \
                createSamples(
                    DAT,
                    j,
                    W_erf[:, 4].flatten(),
                    n_sample,
                    global_log[j],
                    [global_minVal[j], global_maxVal[j]],
                    overshoot,
                    ERF_flag,
                )

        else:

            C_erf = C_EW

        samples_erf[:, j] = C_erf

    if (normalizeSum):

        for triplet in globalSum:

            # print(triplet[0],triplet[1],triplet[2])

            sum_samples = np.sum(samples_EW[:, n_SQ + triplet[0]:n_SQ +
                                            triplet[1] + 1],
                                 axis=1)
            sum_samples = np.expand_dims(sum_samples, axis=1)
            samples_EW[:,
                       n_SQ + triplet[0]:n_SQ + triplet[1] + 1] /= sum_samples
            samples_EW[:,
                       n_SQ + triplet[0]:n_SQ + triplet[1] + 1] *= triplet[2]

            sum_samples = np.sum(samples[:, n_SQ + triplet[0]:n_SQ +
                                         triplet[1] + 1],
                                 axis=1)
            sum_samples = np.expand_dims(sum_samples, axis=1)
            samples[:, n_SQ + triplet[0]:n_SQ + triplet[1] + 1] /= sum_samples
            samples[:, n_SQ + triplet[0]:n_SQ + triplet[1] + 1] *= triplet[2]

            sum_samples = np.sum(samples_erf[:, n_SQ + triplet[0]:n_SQ +
                                             triplet[1] + 1],
                                 axis=1)
            sum_samples = np.expand_dims(sum_samples, axis=1)
            samples_erf[:,
                        n_SQ + triplet[0]:n_SQ + triplet[1] + 1] /= sum_samples
            samples_erf[:,
                        n_SQ + triplet[0]:n_SQ + triplet[1] + 1] *= triplet[2]

    print("       Computing quantiles and mean")

    if verbose:
        print(" j   quan05    quan50     qmean    quan95")

    for j in np.arange(n_SQ + n_TQ):

        if global_log[j]:

            qmean_EW = 10.0**np.mean(np.log10(samples_EW[:, j]))

        else:

            qmean_EW = np.mean(samples_EW[:, j])

        quan05_EW = np.quantile(samples_EW[:, j], 0.05)
        quan50_EW = np.quantile(samples_EW[:, j], 0.5)
        quan95_EW = np.quantile(samples_EW[:, j], 0.95)

        if verbose:
            print(label_indexes[j] + " %.2E %.2E %.2E %.2E" %
                  (quan05_EW, quan50_EW, qmean_EW, quan95_EW))

        q_EW[j, 0] = quan05_EW
        q_EW[j, 1] = quan50_EW
        q_EW[j, 2] = quan95_EW
        q_EW[j, 3] = qmean_EW

        if Cooke_flag != 0:

            if global_log[j]:

                qmean = 10.0**np.mean(np.log10(samples[:, j]))

            else:

                qmean = np.mean(samples[:, j])

            quan05 = np.quantile(samples[:, j], 0.05)
            quan50 = np.quantile(samples[:, j], 0.5)
            quan95 = np.quantile(samples[:, j], 0.95)

            if verbose:
                print(label_indexes[j] + " %.2E %.2E %.2E %.2E" %
                      (quan05, quan50, qmean, quan95))

            q_Cooke[j, 0] = quan05
            q_Cooke[j, 1] = quan50
            q_Cooke[j, 2] = quan95
            q_Cooke[j, 3] = qmean

        if ERF_flag > 0:

            if global_log[j]:

                qmean_erf = 10.0**np.mean(np.log10(samples_erf[:, j]))

            else:

                qmean_erf = np.mean(samples_erf[:, j])

            quan05_erf = np.quantile(samples_erf[:, j], 0.05)
            quan50_erf = np.quantile(samples_erf[:, j], 0.5)
            quan95_erf = np.quantile(samples_erf[:, j], 0.95)

            if verbose:
                print(label_indexes[j] + " %.2E %.2E %.2E %.2E" %
                      (quan05_erf, quan50_erf, qmean_erf, quan95_erf))

            q_erf[j, 0] = quan05_erf
            q_erf[j, 1] = quan50_erf
            q_erf[j, 2] = quan95_erf
            q_erf[j, 3] = qmean_erf

    return q_Cooke, q_erf, q_EW, samples, samples_erf, samples_EW


def main(argv):

    verbose = False

    current_path = os.getcwd()

    if isinstance(argv, str):

        print(argv)
        repository = argv
        path = current_path + "/ELICITATIONS/" + repository

        # Check whether the specified path exists or not
        repoExists = os.path.exists(path)

    else:

        repoExists = False

    if repoExists:

        print("")
        print(repository + " found")
        os.chdir(path)

    else:

        elicitation_list = next(os.walk(current_path + "/ELICITATIONS/"))
        print("List of elicitations:")

        if len(elicitation_list[1]) == 1:

            userinput = 0
            cond = False

        else:

            for count, elicitation in enumerate(elicitation_list[1]):

                print(count, elicitation)

            print("")
            cond = True

        while cond:

            userinput = int(input("Enter Elicitation Number: "))
            cond = not (userinput in range(len(elicitation_list[1])))
            if cond:
                print("Integer between 0 and ", len(elicitation_list[1]) - 1)

        repository = elicitation_list[1][userinput]
        if verbose:
            print("repository", repository)
        path = current_path + "/ELICITATIONS/" + repository

    print('')

    os.chdir(path)
    if verbose:
        print("Path: ", path)
    sys.path.append(path)

    from ElicipyDict import output_dir
    from ElicipyDict import datarepo
    from ElicipyDict import target
    from ElicipyDict import elicitation_name
    from ElicipyDict import analysis
    from ElicipyDict import alpha
    from ElicipyDict import overshoot
    from ElicipyDict import cal_power
    from ElicipyDict import EW_flag
    from ElicipyDict import ERF_flag
    from ElicipyDict import Cooke_flag
    from ElicipyDict import n_sample
    from ElicipyDict import postprocessing
    from ElicipyDict import hist_type
    from ElicipyDict import n_bins

    if (Cooke_flag > 0 or ERF_flag > 0):

        seed = True

    else:

        seed = False

    try:

        from ElicipyDict import nolabel_flag

    except ImportError:

        nolabel_flag = False

    try:

        from ElicipyDict import group_list

    except ImportError:

        group_list = [0]

    if len(group_list) > 1 and not (0 in group_list):

        group_list.append(0)

    n_pctl = 3

    # download the data from github repository
    if datarepo == "github":

        from ElicipyDict import user
        from ElicipyDict import github_token
        from ElicipyDict import RepositoryData

        saveDataFromGithub(RepositoryData, user, github_token)

    sys.path.insert(0, os.getcwd())
    input_dir = "DATA"
    csv_file = "questionnaire.csv"

    # change to full path
    output_dir = path + "/" + output_dir
    input_dir = path + "/" + input_dir

    if Cooke_flag < 0:

        try:

            from ElicipyDict import weights_file

            weights_file = input_dir + "/" + weights_file
            # Check whether the specified file exists or not
            isExist = os.path.exists(weights_file)

            if not isExist:

                print("weights_file does not exist in ", input_dir)
                sys.exit()

        except ImportError:

            print("Please define weights_file")
            sys.exit()

    else:

        weights_file = ''

    # Check whether the specified output path exists or not
    isExist = os.path.exists(output_dir)

    if not isExist:

        # Create a new directory because it does not exist
        os.makedirs(output_dir)
        print("The new directory " + output_dir + " is created!")

    # Read the questionnaire csv file
    (df_indexes_SQ, df_indexes_TQ, SQ_scale, SQ_realization, TQ_scale,
     SQ_minVals, SQ_maxVals, TQ_minVals, TQ_maxVals, SQ_units, TQ_units,
     SQ_LongQuestion, TQ_LongQuestion, SQ_question, TQ_question, idx_list,
     global_scale, global_log, label_indexes, parents, global_idxMin,
     global_idxMax, global_sum50,
     df_quest) = read_questionnaire(input_dir, csv_file, seed, target)

    # Read an renumber index_groups
    try:

        from ElicipyDict import index_groups

    except ImportError:

        index_groups = []

    df_quest.reset_index(inplace=True, drop=True)

    for count, group in enumerate(index_groups):

        new_indexes_group = []

        for idx in group:

            indices = df_quest.index[
                (df_quest["IDX"] == idx)
                & (df_quest.QUEST_TYPE.str.contains("target"))]

            new_indexes_group.append(int(indices[0]))

        index_groups[count] = new_indexes_group

    # Read an renumber trend_groups
    try:

        from ElicipyDict import trend_groups

    except ImportError:

        trend_groups = []

    df_quest.reset_index(inplace=True, drop=True)

    for count, group in enumerate(trend_groups):

        new_trend_group = []

        for idx in group:

            indices = df_quest.index[
                (df_quest["IDX"] == idx)
                & (df_quest.QUEST_TYPE.str.contains("target"))]

            new_trend_group.append(int(indices[0]))

        trend_groups[count] = new_trend_group

    # Read an renumber violin_groups
    try:

        from ElicipyDict import violin_groups

    except ImportError:

        violin_groups = []

    df_quest.reset_index(inplace=True, drop=True)

    for count, group in enumerate(violin_groups):

        new_violin_group = []

        for idx in group:

            indices = df_quest.index[
                (df_quest["IDX"] == idx)
                & (df_quest.QUEST_TYPE.str.contains("target"))]

            new_violin_group.append(int(indices[0]))

        violin_groups[count] = new_violin_group

    # Read an renumber pie_groups
    try:

        from ElicipyDict import pie_groups

    except ImportError:

        pie_groups = []

    df_quest.reset_index(inplace=True, drop=True)

    for count, group in enumerate(pie_groups):

        new_pie_group = []

        for idx in group:

            indices = df_quest.index[
                (df_quest["IDX"] == idx)
                & (df_quest.QUEST_TYPE.str.contains("target"))]

            new_pie_group.append(int(indices[0]))

        pie_groups[count] = new_pie_group

    try:

        from ElicipyDict import normalizeSum

    except ImportError:

        normalizeSum = False

    globalSum_temp = zip(global_idxMin, global_idxMax, global_sum50)

    globalSum_temp = list(set(globalSum_temp))

    globalSum = []
    for triplet in globalSum_temp:

        try:

            int(triplet[0])
            flag = True

        except ValueError:

            flag = False

        # print(flag)
        if flag:

            index_first = idx_list.index(int(triplet[0]))
            index_last = idx_list.index(int(triplet[1]))
            globalSum.append([index_first, index_last, triplet[2]])

    # Read the asnwers of all the experts
    group = 0
    write_flag = True

    print("STEP2: Reading all answers to define ranges")
    n_experts, n_SQ, n_TQ, SQ_array, TQ_array, NS_experts = read_answers(
        input_dir, csv_file, group, n_pctl, df_indexes_SQ, df_indexes_TQ, seed,
        target, output_dir, elicitation_name, write_flag, label_indexes)

    save_dtt_rll(input_dir, n_experts, n_SQ, n_TQ, df_quest, target,
                 SQ_realization, SQ_scale, SQ_array, TQ_scale, TQ_array)

    # alpha_nominal, alpha_interval = /
    # calculate_alpha(TQ_array, overshoot,TQ_scale)

    weight = np.ones(n_experts)
    indexMean, indexStd, indexQuantiles = calculate_index(
        TQ_array, weight, TQ_scale)

    # for i in range(n_TQ):

    # print("TQ " + label_indexes[i+n_SQ],indexMean[i],indexStd[i], /
    # indexQuantiles[i,:])
    # string = "%.2E, " % indexMean[i] + "%.2E, " % indexStd[i] + "Target"
    # print(string)
    # print(ciao)

    minval_all = np.zeros(n_SQ + n_TQ)
    minval_all[0:n_SQ] = np.amin(SQ_array[:, 0, :], axis=0)
    minval_all[n_SQ:] = np.amin(TQ_array[:, 0, :], axis=0)

    maxval_all = np.zeros(n_SQ + n_TQ)
    maxval_all[0:n_SQ] = np.amax(SQ_array[:, 2, :], axis=0)
    maxval_all[n_SQ:] = np.amax(TQ_array[:, 2, :], axis=0)

    global_units = SQ_units + TQ_units

    global_minVal = SQ_minVals + TQ_minVals
    global_maxVal = SQ_maxVals + TQ_maxVals
    global_units = SQ_units + TQ_units
    global_longQuestion = SQ_LongQuestion + TQ_LongQuestion
    global_shortQuestion = SQ_question + TQ_question

    if verbose:
        print("")
        print("Answer ranges")

    if verbose:
        print(df_indexes_SQ, df_indexes_TQ)

    for i in range(n_SQ + n_TQ):

        if global_log[i]:

            deltalogval_all = np.log10(maxval_all[i]) - np.log10(minval_all[i])

            minval_all[i] = 10**(np.log10(minval_all[i]) -
                                 0.1 * deltalogval_all)
            maxval_all[i] = 10**(np.log10(maxval_all[i]) +
                                 0.1 * deltalogval_all)

        else:

            deltaval_all = maxval_all[i] - minval_all[i]
            minval_all[i] = minval_all[i] - 0.1 * deltaval_all
            maxval_all[i] = maxval_all[i] + 0.1 * deltaval_all

            if global_units[i] == "%":

                minval_all[i] = np.maximum(minval_all[i], 0)
                maxval_all[i] = np.minimum(maxval_all[i], 100)

        minval_all[i] = np.maximum(minval_all[i], global_minVal[i])
        maxval_all[i] = np.minimum(maxval_all[i], global_maxVal[i])

        if verbose:
            print(i, minval_all[i], maxval_all[i])

    q_Cooke = np.zeros((len(global_scale), 4))
    q_erf = np.zeros((len(global_scale), 4))
    q_EW = np.zeros((len(global_scale), 4))

    if len(group_list) > 1:

        q_Cooke_groups = np.zeros((len(global_scale), 4, len(group_list)))
        q_erf_groups = np.zeros((len(global_scale), 4, len(group_list)))
        q_EW_groups = np.zeros((len(global_scale), 4, len(group_list)))

    for count, group in enumerate(group_list):

        # Read the asnwers of the experts
        write_flag = False

        print("STEP" + str(3 + 2 * count) + ": Reading answers for group ",
              count)

        n_experts, n_SQ, n_TQ, SQ_array, TQ_array, NS_experts = read_answers(
            input_dir, csv_file, group, n_pctl, df_indexes_SQ, df_indexes_TQ,
            seed, target, output_dir, elicitation_name, write_flag,
            label_indexes)

        if not target:

            # if we do not read the target questions, set empty array
            n_TQ = 0
            TQ_array = np.zeros((n_experts, n_pctl, n_TQ))

        if seed:

            realization = np.zeros(TQ_array.shape[2] + SQ_array.shape[2])
            realization[0:SQ_array.shape[2]] = SQ_realization

            if verbose:

                print("")
                print("Realization")
                print(realization)
                print()

        else:

            realization = []

        if analysis:

            tree = {"IDX": idx_list, "SHORT_Q": TQ_question}
            df_tree = pd.DataFrame(data=tree)

            # ----------------------------------------- #
            # ------------ Compute weights ------------ #
            # ----------------------------------------- #

            if group == 0:

                alpha_analysis = alpha

            else:

                alpha_analysis = 0.0

            W, W_erf, Weqok, W_gt0, Werf_gt0, expin = answer_analysis(
                input_dir, csv_file, n_experts, n_SQ, n_TQ, SQ_array, TQ_array,
                realization, global_scale, global_log, alpha_analysis,
                overshoot, cal_power, ERF_flag, Cooke_flag, seed, NS_experts,
                weights_file)

            # ----------------------------------------- #
            # ------ Create samples and bar plots ----- #
            # ----------------------------------------- #

            q_Cooke, q_erf, q_EW, samples, samples_erf, samples_EW = \
                create_samples(group, n_experts, n_SQ, n_TQ, n_pctl, SQ_array,
                               TQ_array, n_sample, W, W_erf, Weqok, W_gt0,
                               Werf_gt0, expin, global_log, global_minVal,
                               global_maxVal, label_indexes, ERF_flag,
                               abs(Cooke_flag), EW_flag, overshoot, globalSum,
                               normalizeSum)

            try:

                from ElicipyDict import delta_ratio_flag

            except ImportError:

                delta_ratio_flag = False

            if delta_ratio_flag:

                delta_ratio_Cooke = np.zeros(n_TQ)
                delta_ratio_erf = np.zeros(n_TQ)
                delta_ratio_EW = np.zeros(n_TQ)

                for j in range(n_TQ):

                    delta_perc = np.zeros(n_experts)

                    if global_log[n_SQ + j]:

                        delta_perc[:] = np.log10(TQ_array[:, 2, j]) - np.log10(
                            TQ_array[:, 0, j])

                    else:

                        delta_perc[:] = TQ_array[:, 2, j] - TQ_array[:, 0, j]

                    delta_perc_mean = np.mean(delta_perc, axis=0)

                    if Cooke_flag != 0:

                        if global_log[n_SQ + j]:

                            delta_q_Cooke = np.log10(
                                q_Cooke[n_SQ + j, 2]) - np.log10(
                                    q_Cooke[n_SQ + j, 0])

                        else:

                            delta_q_Cooke = q_Cooke[n_SQ + j,
                                                    2] - q_Cooke[n_SQ + j, 0]

                        delta_ratio_Cooke[j] = delta_perc_mean / delta_q_Cooke

                    if ERF_flag > 0:

                        if global_log[n_SQ + j]:

                            delta_q_erf = np.log10(
                                q_erf[n_SQ + j, 2]) - np.log10(q_erf[n_SQ + j,
                                                                     0])

                        else:

                            delta_q_erf = q_erf[n_SQ + j, 2] - q_erf[n_SQ + j,
                                                                     0]

                        delta_ratio_erf[j] = delta_perc_mean / delta_q_erf

                    if EW_flag > 0:

                        if global_log[n_SQ + j]:

                            delta_q_EW = np.log10(
                                q_EW[n_SQ + j, 2]) - np.log10(q_EW[n_SQ + j,
                                                                   0])

                        else:

                            delta_q_EW = q_EW[n_SQ + j, 2] - q_EW[n_SQ + j, 0]

                        delta_ratio_EW[j] = delta_perc_mean / delta_q_EW

            print(
                "STEP" + str(4 + 2 * count) +
                ": Creating distribution plots for group ", count)

            create_barplot(group, n_SQ, n_TQ, n_sample, global_log,
                           global_minVal, global_maxVal, global_units,
                           TQ_units,
                           label_indexes, minval_all, maxval_all, ERF_flag,
                           abs(Cooke_flag), EW_flag, hist_type, output_dir,
                           elicitation_name, n_bins, q_Cooke, q_erf, q_EW,
                           samples, samples_erf, samples_EW)

            if len(group_list) > 1:

                q_Cooke_groups[:, :, count] = q_Cooke
                q_erf_groups[:, :, count] = q_erf
                q_EW_groups[:, :, count] = q_EW

    # ----------------------------------------- #
    # ---------- Save samples on csv ---------- #
    # ----------------------------------------- #

    if analysis:

        csv_name = output_dir + "/" + elicitation_name + "_weights.csv"

        Weqok_100 = [100.0 * elem for elem in Weqok]

        Weqok_formatted = ["%.2f" % elem for elem in Weqok_100]

        d = {"index": range(1, n_experts + 1), "Weq": Weqok_formatted}
        df = pd.DataFrame(data=d)

        if Cooke_flag != 0:

            df.insert(loc=2, column="WCooke", value=W_gt0)

        if ERF_flag > 0:

            df.insert(loc=2, column="WERF", value=Werf_gt0)

        df.to_csv(csv_name, index=False)

    if analysis and target:

        targets = ["target_" + str(i + 1).zfill(2) for i in range(n_TQ)]

        df_tree["EW_5"] = q_EW[n_SQ:, 0]
        df_tree["EW_50"] = q_EW[n_SQ:, 1]
        df_tree["EW_95"] = q_EW[n_SQ:, 2]
        df_tree["EW_MEAN"] = q_EW[n_SQ:, 3]

        if Cooke_flag != 0:

            df_tree["COOKE_5"] = q_Cooke[n_SQ:, 0]
            df_tree["COOKE_50"] = q_Cooke[n_SQ:, 1]
            df_tree["COOKE_95"] = q_Cooke[n_SQ:, 2]
            df_tree["COOKE_MEAN"] = q_Cooke[n_SQ:, 3]

            csv_name = output_dir + "/" + elicitation_name + "_samples.csv"
            np.savetxt(
                csv_name,
                samples[:, n_SQ:],
                header=",".join(targets),
                comments="",
                delimiter=",",
                fmt="%1.4e",
            )

        if ERF_flag > 0:

            df_tree["ERF_5"] = q_erf[n_SQ:, 0]
            df_tree["ERF_50"] = q_erf[n_SQ:, 1]
            df_tree["ERF_95"] = q_erf[n_SQ:, 2]
            df_tree["ERF_MEAN"] = q_erf[n_SQ:, 3]

            csv_name = output_dir + "/" + elicitation_name + "_samples_erf.csv"
            np.savetxt(
                csv_name,
                samples_erf[:, n_SQ:],
                header=",".join(targets),
                comments="",
                delimiter=",",
                fmt="%1.4e",
            )

        if EW_flag > 0:

            csv_name = output_dir + "/" + elicitation_name + "_samples_EW.csv"
            np.savetxt(
                csv_name,
                samples_EW[:, n_SQ:],
                header=",".join(targets),
                comments="",
                delimiter=",",
                fmt="%1.4e",
            )

        df_tree["PARENT"] = parents
        df_tree.to_csv("tree.csv", index=False)

    if not postprocessing:

        print("Analysis completed!")
        sys.exit()

    counter_plot = 0

    if analysis:

        # ------------------------------------------ #
        # ---------- Create index figures ---------- #
        # ------------------------------------------ #

        indexMean_EW, indexStd_EW, indexQuantiles_EW = calculate_index(
            TQ_array, Weqok, TQ_scale)

        if Cooke_flag != 0:

            indexMean_Cooke, indexStd_Cooke, indexQuantiles_Cooke = \
                calculate_index(TQ_array, W_gt0, TQ_scale)

        else:

            indexMean_Cooke = indexMean_EW
            indexStd_Cooke = indexStd_EW
            indexQuantiles_Cooke = indexQuantiles_EW

        if ERF_flag > 0:

            indexMean_erf, indexStd_erf, indexQuantiles_erf = calculate_index(
                TQ_array, Werf_gt0, TQ_scale)

        else:

            indexMean_erf = indexMean_EW
            indexStd_erf = indexStd_EW
            indexQuantiles_erf = indexQuantiles_EW

        if len(index_groups) > 0:

            print("STEP" + str(3 + counter_plot + 2 * len(group_list)) +
                  ": Creating index plots")
            counter_plot += 1

        for count, index_group in enumerate(index_groups):

            create_figure_index(
                count,
                index_group,
                n_SQ,
                label_indexes,
                indexMean_EW,
                indexStd_EW,
                indexMean_Cooke,
                indexStd_Cooke,
                indexMean_erf,
                indexStd_erf,
                indexQuantiles_EW,
                indexQuantiles_Cooke,
                indexQuantiles_erf,
                global_units,
                abs(Cooke_flag),
                ERF_flag,
                EW_flag,
                global_log,
                output_dir,
                elicitation_name,
            )

        # ------------------------------------------ #
        # --------- Create trend. figures ---------- #
        # ------------------------------------------ #

        if len(trend_groups) > 0:

            print("STEP" + str(3 + counter_plot + 2 * len(group_list)) +
                  ": Creating trend plots")
            counter_plot += 1

        for count, trend_group in enumerate(trend_groups):

            create_figure_trend(
                count,
                trend_group,
                n_SQ,
                label_indexes,
                q_EW,
                q_Cooke,
                q_erf,
                global_units,
                abs(Cooke_flag),
                ERF_flag,
                EW_flag,
                global_log,
                TQ_minVals,
                TQ_maxVals,
                output_dir,
                elicitation_name,
            )

        # ------------------------------------------ #
        # --------- Create violin figures ---------- #
        # ------------------------------------------ #

        if len(violin_groups) > 0:

            print("STEP" + str(3 + counter_plot + 2 * len(group_list)) +
                  ": Creating violin plots")
            counter_plot += 1

        for count, violin_group in enumerate(violin_groups):

            create_figure_violin(
                count,
                violin_group,
                n_SQ,
                label_indexes,
                samples_EW,
                samples,
                samples_erf,
                q_EW,
                q_Cooke,
                q_erf,
                global_units,
                abs(Cooke_flag),
                ERF_flag,
                EW_flag,
                global_log,
                TQ_minVals,
                TQ_maxVals,
                output_dir,
                elicitation_name,
            )

        # ------------------------------------------ #
        # ----------- Create pie figures ----------- #
        # ------------------------------------------ #

        if len(pie_groups) > 0:

            print("STEP" + str(3 + counter_plot + 2 * len(group_list)) +
                  ": Creating pie charts")
            counter_plot += 1

        for count, pie_group in enumerate(pie_groups):

            create_figure_pie(count, pie_group, n_SQ,
                              label_indexes, q_EW, q_Cooke, q_erf,
                              abs(Cooke_flag), ERF_flag, EW_flag, output_dir,
                              elicitation_name)

    # ----------------------------------------- #
    # --------- Create answ. figures ---------- #
    # ----------------------------------------- #

    n_panels = int(np.ceil(n_experts / max_len_plot))

    print("STEP" + str(3 + counter_plot + 2 * len(group_list)) +
          ": Creating answer plots")

    for h in np.arange(n_SQ + n_TQ):

        if (n_SQ + n_TQ > 1):

            printProgressBar(h, n_SQ + n_TQ - 1, prefix='      ')

        for k in np.arange(n_panels):

            create_figure_answers(h, k, n_experts, max_len_plot, n_SQ,
                                  SQ_array, TQ_array, realization, analysis,
                                  abs(Cooke_flag), ERF_flag, EW_flag,
                                  global_units, output_dir, q_Cooke, q_erf,
                                  q_EW, elicitation_name, global_log,
                                  label_indexes, nolabel_flag)

    # ----------------------------------------- #
    # ------- Create .pptx presentation ------- #
    # ----------------------------------------- #

    print("STEP" + str(4 + counter_plot + 2 * len(group_list)) +
          ": Creating presentation")

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    left = Inches(2)
    top = Inches(1.5)

    # ------------- Title slide ----------------#
    lyt = prs.slide_layouts[0]  # choosing a slide layout
    slide = prs.slides.add_slide(lyt)  # adding a slide
    title = slide.shapes.title  # assigning a title
    subtitle = slide.placeholders[1]  # placeholder for subtitle
    title.text = "Expert elicitation - " + elicitation_name  # title

    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.name = "Helvetica"

    Current_Date_Formatted = datetime.datetime.today().strftime("%d-%b-%Y")

    subtitle.text = Current_Date_Formatted  # subtitle

    subtitle_para = slide.shapes.placeholders[1].text_frame.paragraphs[0]
    subtitle_para.font.name = "Helvetica"

    logofile = current_path + "/logo.png"

    slide.shapes.add_picture(logofile,
                             left + Inches(11.3),
                             top + Inches(5.4),
                             width=Inches(2.4))

    title_slide_layout = prs.slide_layouts[5]

    # ------------- Weights slide -------------#

    if analysis and seed:

        slide = prs.slides.add_slide(title_slide_layout)

        text_title = "Experts' weights"

        title_shape = slide.shapes.title
        title_shape.text = text_title
        title_shape.top = Inches(3.0)
        title_shape.width = Inches(15)
        title_shape.height = Inches(2)
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.name = "Helvetica"
        title_para.font.size = Pt(54)
        add_date(slide)
        add_small_logo(slide, left, top, logofile)

        n_tables = int(np.ceil(len(W_gt0) / max_len_table))

        for i_table in range(n_tables):

            slide = prs.slides.add_slide(title_slide_layout)

            text_title = "Experts' weights"
            add_title(slide, text_title)

            # ---add table weights to slide---
            x, y, cx = Inches(2), Inches(1.7), Inches(8)

            fisrt_j = i_table * max_len_table
            last_j = np.minimum((i_table + 1) * max_len_table, len(W_gt0))

            if ERF_flag > 0:

                shape = slide.shapes.add_table(last_j - fisrt_j + 1, 3, x, y,
                                               cx,
                                               MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)

            else:

                shape = slide.shapes.add_table(last_j - fisrt_j + 1, 2, x, y,
                                               cx,
                                               MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)

            table = shape.table

            cell = table.cell(0, 0)
            cell.text = "Expert ID"

            cell = table.cell(0, 1)
            cell.text = "Cooke"

            if ERF_flag > 0:

                cell = table.cell(0, 2)
                cell.text = "ERF"

            for j in np.arange(fisrt_j, last_j):
                j_mod = np.remainder(j, max_len_table)
                cell = table.cell(j_mod + 1, 0)
                cell.text = "Exp" + str(expin[j])

                cell = table.cell(j_mod + 1, 1)

                if W_gt0[j] > 0.0:

                    cell.text = "%6.2f" % W_gt0[j]

                else:

                    cell.text = "Below threshold"

                if ERF_flag > 0:

                    cell = table.cell(j_mod + 1, 2)
                    cell.text = "%6.2f" % Werf_gt0[j]

            for cell in iter_cells(table):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(12)

            if EW_flag:

                text_box = "Equal weight = " + f"{Weqok[0]*100:.2f}"
                add_text_box(slide, Inches(12), Inches(2), text_box, 18)

                text_box = ("For Cookes' method, below threshold weight does "
                            "not mean zero score. It simply means that this "
                            "expert's knowledge was already contributed by "
                            "other experts and adding this expert would not "
                            "change significantly the results.")
                add_text_box(slide, Inches(12), Inches(3), text_box, 18)

            add_date(slide)
            add_small_logo(slide, left, top, logofile)

    # ------------- Answers slides ------------#

    if seed:

        slide = prs.slides.add_slide(title_slide_layout)

        text_title = "Seed Answers"

        title_shape = slide.shapes.title
        title_shape.text = text_title
        title_shape.top = Inches(3.0)
        title_shape.width = Inches(15)
        title_shape.height = Inches(2)
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.name = "Helvetica"
        title_para.font.size = Pt(54)
        add_date(slide)
        add_small_logo(slide, left, top, logofile)

    for h in np.arange(n_SQ + n_TQ):

        if h == n_SQ:

            slide = prs.slides.add_slide(title_slide_layout)

            text_title = "Target Answers"

            title_shape = slide.shapes.title
            title_shape.text = text_title
            title_shape.top = Inches(3.0)
            title_shape.width = Inches(15)
            title_shape.height = Inches(2)
            title_para = slide.shapes.title.text_frame.paragraphs[0]
            title_para.font.name = "Helvetica"
            title_para.font.size = Pt(54)
            add_date(slide)
            add_small_logo(slide, left, top, logofile)

        text_box = global_longQuestion[h]

        if h >= n_SQ:

            j = h - n_SQ
            string = "Target"

            if len(text_box) < 250:

                fontsize = 18

            else:

                fontsize = 18.0 * np.sqrt(250.0 / len(text_box))

        else:

            j = h
            string = "Seed"

            if len(text_box) < 500:

                fontsize = 18

            else:

                fontsize = 18.0 * np.sqrt(500.0 / len(text_box))

        for k in np.arange(n_panels):

            slide = prs.slides.add_slide(title_slide_layout)

            add_text_box(slide, left, top, text_box, fontsize)

            figname = (output_dir + "/" + "Itemwise_PNGPDF" + "/" +
                       elicitation_name + "_" + string + "_" +
                       str(j + 1).zfill(2) + "_" + str(k + 1).zfill(2) +
                       ".png")
            add_figure(slide, figname, left, top, Inches(10))

            if analysis and (string == "Target"):

                figname = (output_dir + "/" + "Groups_PNGPDF" + "/" +
                           elicitation_name + "_cum_group0_" +
                           str(j + 1).zfill(2) + ".png")

                width = Inches(2.5)
                add_small_figure(slide, figname, left + Inches(1.1),
                                 top + Inches(4.9), width)

                figname = (output_dir + "/" + "Groups_PNGPDF" + "/" +
                           elicitation_name + "_PDF_group0_" +
                           str(j + 1).zfill(2) + ".png")

                width = Inches(2.5)
                add_small_figure(slide, figname, left - Inches(1.2),
                                 top + Inches(4.9), width)

                if global_idxMin[h] != global_idxMax[h]:

                    longQ_NB = ("N.B. The sum of 50%iles for questions " +
                                str(global_idxMin[h]) + "-" +
                                str(global_idxMax[h]) + " have to sum to " +
                                str(global_sum50[h]) + ".")

                    add_text_box(slide, left, top + Inches(3.79), longQ_NB, 13)

            add_date(slide)
            add_small_logo(slide, left, top, logofile)

            text_title = global_shortQuestion[h]
            add_title(slide, text_title)

    # ------------- Pctls slides -------------#

    if analysis and target:

        slide = prs.slides.add_slide(title_slide_layout)

        text_title = "Target Percentiles"

        title_shape = slide.shapes.title
        title_shape.text = text_title
        title_shape.top = Inches(3.0)
        title_shape.width = Inches(15)
        title_shape.height = Inches(2)
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.name = "Helvetica"
        title_para.font.size = Pt(54)
        add_date(slide)
        add_small_logo(slide, left, top, logofile)

        n_tables = int(np.ceil(n_TQ / max_len_tableB))

        for i_table in range(n_tables):

            slide = prs.slides.add_slide(prs.slide_layouts[5])

            fisrt_j = i_table * max_len_tableB
            last_j = np.minimum((i_table + 1) * max_len_tableB, n_TQ)

            n_columns = 1

            if Cooke_flag != 0:

                n_columns += 4

            if EW_flag > 0:

                n_columns += 4

            if ERF_flag > 0:

                n_columns += 4

            # ---add table to slide---
            x, y, cx = Inches(1), Inches(2), Inches(14)
            shape = slide.shapes.add_table(
                last_j - fisrt_j + 1,
                n_columns,
                x,
                y,
                cx,
                MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
            )
            table = shape.table

            cell = table.cell(0, 1)
            cell.text = "Q05 (EW)"

            cell = table.cell(0, 2)
            cell.text = "Q50 (EW)"

            cell = table.cell(0, 3)
            cell.text = "Q95 (EW)"

            cell = table.cell(0, 4)
            cell.text = "Qmean (EW)"

            j_column = 4

            if Cooke_flag != 0:

                cell = table.cell(0, j_column + 1)
                cell.text = "Q05 (Cooke)"

                cell = table.cell(0, j_column + 2)
                cell.text = "Q50 (Cooke)"

                cell = table.cell(0, j_column + 3)
                cell.text = "Q95 (Cooke)"

                cell = table.cell(0, j_column + 4)
                cell.text = "Qmean (Cooke)"
                j_column += 4

            if ERF_flag > 0:

                cell = table.cell(0, j_column + 1)
                cell.text = "Q05 (ERF)"

                cell = table.cell(0, j_column + 2)
                cell.text = "Q50 (ERF)"

                cell = table.cell(0, j_column + 3)
                cell.text = "Q95 (ERF)"

                cell = table.cell(0, j_column + 4)
                cell.text = "Qmean (ERF)"

            for h in np.arange(fisrt_j, last_j):

                h_mod = np.remainder(h, max_len_tableB)

                j = h + n_SQ

                cell = table.cell(h_mod + 1, 0)
                cell.text = "TQ " + label_indexes[j]

                for li in range(4):

                    cell = table.cell(h_mod + 1, li + 1)
                    if global_units[j] == "%" and global_log[j] == 0:
                        cell.text = "%6.2f" % q_EW[j, li]
                    else:
                        cell.text = "%.2E" % q_EW[j, li]

                    j_column = 4

                    if Cooke_flag != 0:

                        cell = table.cell(h_mod + 1, j_column + li + 1)
                        if global_units[j] == "%" and global_log[j] == 0:
                            cell.text = "%6.2f" % q_Cooke[j, li]
                        else:
                            cell.text = "%.2E" % q_Cooke[j, li]

                        j_column += 4

                    if ERF_flag > 0:

                        cell = table.cell(h_mod + 1, j_column + li + 1)
                        if global_units[j] == "%" and global_log[j] == 0:
                            cell.text = "%6.2f" % q_erf[j, li]
                        else:
                            cell.text = "%.2E" % q_erf[j, li]

            for cell in iter_cells(table):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)

            add_date(slide)
            add_small_logo(slide, left, top, logofile)

            text_title = "Percentiles of target questions"
            add_title(slide, text_title)

    # ------------- Delta ratio slides -------------#

    if analysis and target and delta_ratio_flag:

        slide = prs.slides.add_slide(title_slide_layout)

        text_title = "Target TEST ANALYSIS"

        title_shape = slide.shapes.title
        title_shape.text = text_title
        title_shape.top = Inches(3.0)
        title_shape.width = Inches(15)
        title_shape.height = Inches(2)
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.name = "Helvetica"
        title_para.font.size = Pt(54)
        add_date(slide)
        add_small_logo(slide, left, top, logofile)

        n_tables = int(np.ceil(n_TQ / max_len_tableB))

        for i_table in range(n_tables):

            slide = prs.slides.add_slide(prs.slide_layouts[5])

            fisrt_j = i_table * max_len_tableB
            last_j = np.minimum((i_table + 1) * max_len_tableB, n_TQ)

            n_columns = 1

            if Cooke_flag != 0:

                n_columns += 1

            if EW_flag > 0:

                n_columns += 1

            if ERF_flag > 0:

                n_columns += 1

            # ---add table to slide---
            x, y, cx = Inches(1), Inches(2), Inches(14)
            shape = slide.shapes.add_table(
                last_j - fisrt_j + 1,
                n_columns,
                x,
                y,
                cx,
                MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
            )
            table = shape.table

            cell = table.cell(0, 1)
            cell.text = "Delta ratio (EW)"

            j_column = 1

            if Cooke_flag != 0:

                cell = table.cell(0, j_column + 1)
                cell.text = "Delta ratio  (Cooke)"

                j_column += 1

            if ERF_flag > 0:

                cell = table.cell(0, j_column + 1)
                cell.text = "Delta ratio  (ERF)"

            for h in np.arange(fisrt_j, last_j):

                h_mod = np.remainder(h, max_len_tableB)

                j = h + n_SQ

                cell = table.cell(h_mod + 1, 0)
                cell.text = "TQ " + label_indexes[j]

                cell = table.cell(h_mod + 1, 1)
                cell.text = "%.2E" % delta_ratio_EW[h]

                j_column = 1

                if Cooke_flag != 0:

                    cell = table.cell(h_mod + 1, j_column + 1)
                    cell.text = "%.2E" % delta_ratio_Cooke[h]

                    j_column += 1

                if ERF_flag > 0:

                    cell = table.cell(h_mod + 1, j_column + 1)
                    cell.text = "%.2E" % delta_ratio_erf[h]

            for cell in iter_cells(table):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)

            add_date(slide)
            add_small_logo(slide, left, top, logofile)

            text_title = "Target questions TEST ANALYSIS"
            add_title(slide, text_title)

    # ----------- Index groups slides --------#

    if analysis and len(index_groups) > 0:

        slide = prs.slides.add_slide(title_slide_layout)

        text_title = "Index plots"

        title_shape = slide.shapes.title
        title_shape.text = text_title
        title_shape.top = Inches(3.0)
        title_shape.width = Inches(15)
        title_shape.height = Inches(2)
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.name = "Helvetica"
        title_para.font.size = Pt(54)
        add_date(slide)
        add_small_logo(slide, left, top, logofile)

        for count, index_group in enumerate(index_groups):

            slide = prs.slides.add_slide(title_slide_layout)

            figname = (output_dir + "/" + "Index_PNGPDF" + "/" +
                       elicitation_name + "_index_" + str(count + 1).zfill(2) +
                       ".png")

            text_box = ""
            for i in index_group:

                text_box = (text_box + "TQ" + label_indexes[i] + ". " +
                            TQ_question[i - n_SQ] + ".\n\n")

            if len(index_group) < 6:

                fontsize = 18

            else:

                fontsize = 18.0 * np.sqrt(6.0 / len(index_group))

            add_text_box(slide, left, top, text_box, fontsize)

            add_date(slide)
            add_small_logo(slide, left, top, logofile)
            add_figure(slide, figname, left - Inches(0.8), top, Inches(10))

            text_title = "Target questions Group " + str(count + 1)
            add_title(slide, text_title)

            text_box = "Intervals show mean and standard error of the index."
            fontsize = 13

            add_text_box(slide, left, top+Inches(6.0), text_box, fontsize)

    # ----------- Trend groups slides --------#

    if analysis and len(trend_groups) > 0:

        slide = prs.slides.add_slide(title_slide_layout)

        text_title = "Trend plots"

        title_shape = slide.shapes.title
        title_shape.text = text_title
        title_shape.top = Inches(3.0)
        title_shape.width = Inches(15)
        title_shape.height = Inches(2)
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.name = "Helvetica"
        title_para.font.size = Pt(54)
        add_date(slide)
        add_small_logo(slide, left, top, logofile)

        for count, trend_group in enumerate(trend_groups):

            slide = prs.slides.add_slide(title_slide_layout)

            figname = (output_dir + "/" + "Trend_PNGPDF" + "/" +
                       elicitation_name + "_trend_" + str(count + 1).zfill(2) +
                       ".png")

            text_box = ""
            for i in trend_group:

                text_box = (text_box + "TQ" + label_indexes[i] + ". " +
                            TQ_question[i - n_SQ] + ".\n\n")

            if len(trend_group) < 6:

                fontsize = 18

            else:

                fontsize = 18.0 * np.sqrt(6.0 / len(trend_group))

            add_text_box(slide, left, top, text_box, fontsize)

            add_date(slide)
            add_small_logo(slide, left, top, logofile)
            add_figure(slide, figname, left - Inches(0.8), top, Inches(10))

            text_title = "Target questions Group " + str(count + 1)
            add_title(slide, text_title)

    # ----------- Violin groups slides --------#

    if analysis and len(violin_groups) > 0:

        slide = prs.slides.add_slide(title_slide_layout)

        text_title = "Violin plots"

        title_shape = slide.shapes.title
        title_shape.text = text_title
        title_shape.top = Inches(3.0)
        title_shape.width = Inches(15)
        title_shape.height = Inches(2)
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.name = "Helvetica"
        title_para.font.size = Pt(54)
        add_date(slide)
        add_small_logo(slide, left, top, logofile)

        for count, violin_group in enumerate(violin_groups):

            slide = prs.slides.add_slide(title_slide_layout)

            figname = (output_dir + "/" + "Violin_PNGPDF" + "/" +
                       elicitation_name + "_violin_" +
                       str(count + 1).zfill(2) + ".png")

            text_box = ""
            for i in violin_group:

                text_box = (text_box + "TQ" + label_indexes[i] + ". " +
                            TQ_question[i - n_SQ] + ".\n\n")

            if len(violin_group) < 6:

                fontsize = 18

            else:

                fontsize = 18.0 * np.sqrt(6.0 / len(violin_group))

            add_text_box(slide, left, top, text_box, fontsize)

            add_date(slide)
            add_small_logo(slide, left, top, logofile)
            add_figure(slide, figname, left - Inches(0.34), top + Inches(1.0),
                       Inches(10.85))

            text_title = "Target questions Group " + str(count + 1)
            add_title(slide, text_title)

    # ----------- Pie groups slides --------#

    if analysis and len(pie_groups) > 0:

        slide = prs.slides.add_slide(title_slide_layout)

        text_title = "Pie charts"

        title_shape = slide.shapes.title
        title_shape.text = text_title
        title_shape.top = Inches(3.0)
        title_shape.width = Inches(15)
        title_shape.height = Inches(2)
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.name = "Helvetica"
        title_para.font.size = Pt(54)
        add_date(slide)
        add_small_logo(slide, left, top, logofile)

        for count, pie_group in enumerate(pie_groups):

            slide = prs.slides.add_slide(title_slide_layout)

            figname = (output_dir + "/" + "Piecharts_PNGPDF" + "/" +
                       elicitation_name + "_pie_" + str(count + 1).zfill(2) +
                       ".png")

            text_box = ""
            for i in pie_group:

                text_box = (text_box + "TQ" + label_indexes[i] + ". " +
                            TQ_question[i - n_SQ] + ".\n\n")

            if len(pie_group) < 6:

                fontsize = 18

            else:

                fontsize = 18.0 * np.sqrt(6.0 / len(pie_group))

            add_text_box(slide, left, top, text_box, fontsize)

            add_date(slide)
            add_small_logo(slide, left, top, logofile)
            add_figure(slide, figname, left - Inches(0.34), top + Inches(1.0),
                       Inches(10.85))

            text_title = "Target questions Group " + str(count + 1)
            add_title(slide, text_title)

    # ------------ Barplot slides ------------#

    for j in np.arange(n_SQ + n_TQ):

        if analysis:

            if j == n_SQ:

                slide = prs.slides.add_slide(title_slide_layout)

                text_title = "Target Barplots"

                title_shape = slide.shapes.title
                title_shape.text = text_title
                title_shape.top = Inches(3.0)
                title_shape.width = Inches(15)
                title_shape.height = Inches(2)
                title_para = slide.shapes.title.text_frame.paragraphs[0]
                title_para.font.name = "Helvetica"
                title_para.font.size = Pt(54)
                add_date(slide)
                add_small_logo(slide, left, top, logofile)

            if j >= n_SQ:

                slide = prs.slides.add_slide(title_slide_layout)

                figname = (output_dir + "/" + "Barplots_PNGPDF" + "/" +
                           elicitation_name + "_hist_" +
                           str(j - n_SQ + 1).zfill(2) + ".png")

                text_box = TQ_LongQuestion[j - n_SQ]

                if len(text_box) < 400:

                    fontsize = 18

                else:

                    fontsize = 18.0 * np.sqrt(400.0 / len(text_box))

                add_text_box(slide, left, top, text_box, fontsize)

                add_date(slide)
                add_small_logo(slide, left, top, logofile)
                add_figure(slide, figname, left - Inches(0.8), top, Inches(10))

                text_title = TQ_question[j - n_SQ]
                add_title(slide, text_title)

                # ---add table to slide---

                n_rows = 1

                if Cooke_flag != 0:

                    n_rows += 1

                if EW_flag > 0:

                    n_rows += 1

                if ERF_flag > 0:

                    n_rows += 1
                x, y, cx = Inches(1), Inches(6.5), Inches(4.0)
                shape = slide.shapes.add_table(n_rows, 5, x, y, cx,
                                               MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
                table = shape.table

                cell = table.cell(0, 1)
                cell.text = "Q05"

                cell = table.cell(0, 2)
                cell.text = "Q50"

                cell = table.cell(0, 3)
                cell.text = "Q95"

                cell = table.cell(0, 4)
                cell.text = "Qmean"

                j_row = 0

                if Cooke_flag != 0:

                    cell = table.cell(j_row + 1, 0)
                    cell.text = "Cooke"

                    for li in range(4):

                        cell = table.cell(j_row + 1, li + 1)
                        if global_units[j] == "%" and global_log == 0:
                            cell.text = "%6.2f" % q_Cooke[j, li]
                        else:
                            cell.text = "%.2E" % q_Cooke[j, li]

                    j_row += 1

                if ERF_flag > 0:

                    cell = table.cell(j_row + 1, 0)
                    cell.text = "ERF"

                    for li in range(4):

                        cell = table.cell(j_row + 1, li + 1)
                        if global_units[j] == "%" and global_log == 0:
                            cell.text = "%6.2f" % q_erf[j, li]
                        else:
                            cell.text = "%.2E" % q_erf[j, li]

                    j_row += 1

                if EW_flag > 0:

                    cell = table.cell(j_row + 1, 0)
                    cell.text = "EW"

                    for li in range(4):

                        cell = table.cell(j_row + 1, li + 1)
                        if global_units[j] == "%" and global_log == 0:
                            cell.text = "%6.2f" % q_EW[j, li]
                        else:
                            cell.text = "%.2E" % q_EW[j, li]

                    j_row += 1

                for cell in iter_cells(table):
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(8)

    # ------------ Group slides ------------#

    for j in np.arange(n_SQ + n_TQ):

        if analysis and len(group_list) > 1:

            if j == n_SQ:

                slide = prs.slides.add_slide(title_slide_layout)

                text_title = "Groups PDFs"

                title_shape = slide.shapes.title
                title_shape.text = text_title
                title_shape.top = Inches(3.0)
                title_shape.width = Inches(15)
                title_shape.height = Inches(2)
                title_para = slide.shapes.title.text_frame.paragraphs[0]
                title_para.font.name = "Helvetica"
                title_para.font.size = Pt(54)
                add_date(slide)
                add_small_logo(slide, left, top, logofile)

            if j >= n_SQ:

                slide = prs.slides.add_slide(title_slide_layout)

                for count, group in enumerate(group_list):

                    figname = (output_dir + "/" + "Groups_PNGPDF" + "/" +
                               elicitation_name + "_PDF_group" + str(group) +
                               "_" + str(j - n_SQ + 1).zfill(2) + ".png")

                    if group == 0:

                        width = Inches(6.0)
                        add_small_figure(slide, figname, Inches(5.5),
                                         Inches(2.0), width)

                    else:

                        width = Inches(4.22)
                        add_small_figure(
                            slide,
                            figname,
                            Inches(11.18),
                            Inches(2.2) + count * Inches(3.02),
                            width,
                        )

                text_box = TQ_LongQuestion[j - n_SQ]

                if len(text_box) < 200:

                    fontsize = 18

                else:

                    fontsize = 18.0 * np.sqrt(200.0 / len(text_box))

                add_text_box(slide, left, top, text_box, fontsize)

                add_date(slide)
                add_small_logo(slide, left, top, logofile)

                text_title = TQ_question[j - n_SQ]
                add_title(slide, text_title)

                n_rows = 4

                if EW_flag > 0:

                    n_rows += 3

                if ERF_flag > 0:

                    n_rows += 3

                # x, y, cx = Inches(1), Inches(5.5), Inches(4.0)
                x, y, cx = Inches(0.5), Inches(5.5), Inches(5.0)
                shape = slide.shapes.add_table(n_rows, 5, x, y, cx,
                                               MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)

                table = shape.table

                cell = table.cell(0, 1)
                cell.text = "Q05"

                cell = table.cell(0, 2)
                cell.text = "Q50"

                cell = table.cell(0, 3)
                cell.text = "Q95"

                cell = table.cell(0, 4)
                cell.text = "Qmean"

                j_row = 0

                if Cooke_flag != 0:

                    for count, group in enumerate(group_list):

                        cell = table.cell(j_row + 1 + count, 0)

                        if group == 0:
                            cell.text = "Cooke"

                        else:
                            cell.text = "Cooke " + str(group)

                        for li in range(4):

                            cell = table.cell(j_row + 1 + count, li + 1)

                            if global_units[j] == "%" and global_log == 0:
                                cell.text = "%6.3f" % q_Cooke_groups[j, li,
                                                                     count]
                            else:
                                cell.text = "%.2E" % q_Cooke_groups[j, li,
                                                                    count]

                    j_row += len(group_list)

                if ERF_flag > 0:

                    for count, group in enumerate(group_list):

                        cell = table.cell(j_row + 1 + count, 0)

                        if group == 0:
                            cell.text = "ERF"

                        else:
                            cell.text = "ERF " + str(group)

                        for li in range(4):

                            cell = table.cell(j_row + 1 + count, li + 1)

                            if global_units[j] == "%" and global_log == 0:
                                cell.text = "%6.3f" % q_erf_groups[j, li,
                                                                   count]
                            else:
                                cell.text = "%.2E" % q_erf_groups[j, li, count]

                    j_row += len(group_list)

                if EW_flag > 0:

                    for count, group in enumerate(group_list):

                        cell = table.cell(j_row + 1 + count, 0)

                        if group == 0:
                            cell.text = "EW"

                        else:
                            cell.text = "EW " + str(group)

                        for li in range(4):

                            cell = table.cell(j_row + 1 + count, li + 1)

                            if global_units[j] == "%" and global_log == 0:
                                cell.text = "%6.2f" % q_EW_groups[j, li, count]
                            else:
                                cell.text = "%.2E" % q_EW_groups[j, li, count]

                    j_row += len(group_list)

                for cell in iter_cells(table):
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(8)

    prs.save(output_dir + "/" + elicitation_name + ".pptx")  # saving file
    save_dtt_rll(input_dir, n_experts, n_SQ, n_TQ, df_quest, target,
                 SQ_realization, SQ_scale, SQ_array, TQ_scale, TQ_array)


if __name__ == "__main__":
    main(sys.argv[1:])
